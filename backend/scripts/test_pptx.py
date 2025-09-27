#!/usr/bin/env python3
"""
BD Bible PowerPoint Processing Test Script
Tests the python-pptx installation and basic functionality
"""

import os
import sys

# Add parent directory to path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from PIL import Image
import numpy as np

def test_packages():
    """Test that all packages are working"""
    print("🧪 Testing Python Package Installation")
    print("=" * 40)

    # Test python-pptx
    print("\n📊 Testing python-pptx...")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "BD Bible Test Slide"
        print("✅ Can create PowerPoint presentations")
    except Exception as e:
        print(f"❌ Error with python-pptx: {e}")

    # Test Pillow
    print("\n🖼️ Testing Pillow...")
    try:
        img = Image.new('RGB', (100, 100), color='red')
        img_array = np.array(img)
        print(f"✅ Can create images (shape: {img_array.shape})")
    except Exception as e:
        print(f"❌ Error with Pillow: {e}")

    # Test numpy
    print("\n🔢 Testing numpy...")
    try:
        arr = np.array([1, 2, 3, 4, 5])
        print(f"✅ Can create arrays (mean: {arr.mean()})")
    except Exception as e:
        print(f"❌ Error with numpy: {e}")

    print("\n" + "=" * 40)
    print("✨ All tests completed!")

if __name__ == "__main__":
    test_packages()