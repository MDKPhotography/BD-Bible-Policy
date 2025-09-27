#!/usr/bin/env python3
"""
Fill PowerPoint quad chart template with data while preserving all formatting
"""

import os
import sys
import json
from pathlib import Path
from datetime import datetime
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class QuadChartFiller:
    def __init__(self, template_path):
        """Initialize with template path"""
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        # Load template as base
        self.presentation = None
        self.load_template()

    def load_template(self):
        """Load the template presentation"""
        self.presentation = Presentation(self.template_path)

    def preserve_formatting(self, source_paragraph, target_paragraph, new_text):
        """Preserve formatting from source to target paragraph with new text"""
        # Clear existing runs
        for run in target_paragraph.runs:
            run.text = ''

        # If source has runs, copy their formatting
        if source_paragraph.runs:
            source_run = source_paragraph.runs[0]
            new_run = target_paragraph.runs[0] if target_paragraph.runs else target_paragraph.add_run()

            # Set text
            new_run.text = new_text

            # Copy font properties
            if source_run.font.name:
                new_run.font.name = source_run.font.name
            if source_run.font.size:
                new_run.font.size = source_run.font.size
            if source_run.font.bold is not None:
                new_run.font.bold = source_run.font.bold
            if source_run.font.italic is not None:
                new_run.font.italic = source_run.font.italic
            if source_run.font.underline is not None:
                new_run.font.underline = source_run.font.underline

            # Copy color if exists
            if hasattr(source_run.font, 'color') and source_run.font.color:
                if hasattr(source_run.font.color, 'rgb') and source_run.font.color.rgb:
                    new_run.font.color.rgb = source_run.font.color.rgb
        else:
            # No runs in source, just set text
            if target_paragraph.runs:
                target_paragraph.runs[0].text = new_text
            else:
                target_paragraph.add_run().text = new_text

    def replace_text_in_shape(self, shape, placeholder_text, replacement_text):
        """Replace placeholder text with new text while preserving formatting"""
        if not shape.has_text_frame:
            return False

        text_frame = shape.text_frame
        replaced = False

        for paragraph_idx, paragraph in enumerate(text_frame.paragraphs):
            if placeholder_text.lower() in paragraph.text.lower():
                # Store original formatting
                original_alignment = paragraph.alignment
                original_level = paragraph.level

                # Replace text while preserving run formatting
                for run in paragraph.runs:
                    if placeholder_text.lower() in run.text.lower():
                        # Store original font properties
                        orig_font_name = run.font.name
                        orig_font_size = run.font.size
                        orig_font_bold = run.font.bold
                        orig_font_italic = run.font.italic
                        orig_font_underline = run.font.underline
                        orig_font_color = None
                        if hasattr(run.font, 'color') and run.font.color:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                orig_font_color = run.font.color.rgb

                        # Replace text
                        run.text = run.text.replace(placeholder_text, replacement_text)

                        # Restore font properties
                        if orig_font_name:
                            run.font.name = orig_font_name
                        if orig_font_size:
                            run.font.size = orig_font_size
                        if orig_font_bold is not None:
                            run.font.bold = orig_font_bold
                        if orig_font_italic is not None:
                            run.font.italic = orig_font_italic
                        if orig_font_underline is not None:
                            run.font.underline = orig_font_underline
                        if orig_font_color:
                            run.font.color.rgb = orig_font_color

                        replaced = True

                # Restore paragraph properties
                paragraph.alignment = original_alignment
                paragraph.level = original_level

        return replaced

    def replace_text_in_table(self, table, placeholder_text, replacement_text):
        """Replace text in table cells while preserving formatting"""
        replaced = False
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        if placeholder_text.lower() in paragraph.text.lower():
                            # Store original formatting
                            original_text = paragraph.text
                            original_alignment = paragraph.alignment

                            # Replace in runs to preserve formatting
                            for run in paragraph.runs:
                                if placeholder_text.lower() in run.text.lower():
                                    run.text = run.text.replace(placeholder_text, replacement_text)
                                    replaced = True

                            # Restore alignment
                            if original_alignment:
                                paragraph.alignment = original_alignment

        return replaced

    def fill_quad_chart(self, data_mapping):
        """
        Fill the quad chart with provided data

        data_mapping: dict with placeholder text as keys and replacement text as values
        Example:
        {
            "[Company Name]": "Acme Corp",
            "[Technology Title]": "Advanced Radar System",
            "[Technical POC]": "Dr. Jane Smith",
            "[Date]": "2024-01-15",
            "[Quadrant 1 Title]": "Problem Statement",
            "[Quadrant 1 Content]": "Current radar systems lack...",
            ...
        }
        """
        if not self.presentation:
            self.load_template()

        # Process each slide
        for slide in self.presentation.slides:
            # Process all shapes
            for shape in slide.shapes:
                # Handle regular text shapes
                if shape.has_text_frame:
                    for placeholder, replacement in data_mapping.items():
                        self.replace_text_in_shape(shape, placeholder, replacement)

                # Handle tables
                if shape.has_table:
                    for placeholder, replacement in data_mapping.items():
                        self.replace_text_in_table(shape.table, placeholder, replacement)

                # Handle grouped shapes
                if hasattr(shape, 'shapes'):
                    for sub_shape in shape.shapes:
                        if sub_shape.has_text_frame:
                            for placeholder, replacement in data_mapping.items():
                                self.replace_text_in_shape(sub_shape, placeholder, replacement)

    def save(self, output_path):
        """Save the filled presentation"""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(output_path))
        return str(output_path)

def fill_template_from_json(template_path, data_json, output_path):
    """Convenience function to fill template from JSON data"""
    filler = QuadChartFiller(template_path)

    # Parse JSON if string
    if isinstance(data_json, str):
        data = json.loads(data_json)
    else:
        data = data_json

    # Fill the template
    filler.fill_quad_chart(data)

    # Save the result
    return filler.save(output_path)

def main():
    """Example usage"""
    # Template path
    template_path = Path(__file__).parent.parent.parent / "templates" / "powerpoint" / "space-bd-quad-charts.pptx"

    # Example data mapping - these should match placeholders in your template
    example_data = {
        "[Company Name]": "SpaceTech Innovations",
        "[Technology Title]": "Quantum Radar Detection System",
        "[Date]": datetime.now().strftime("%B %d, %Y"),
        "[Technical POC]": "Dr. Sarah Johnson",
        "[Email]": "sarah.johnson@spacetech.com",
        "[Phone]": "(555) 123-4567",

        # Quadrant content examples
        "[Problem Statement]": "Current space surveillance systems lack the resolution and sensitivity needed to track small debris objects in LEO.",
        "[Solution]": "Quantum-enhanced radar technology providing 10x improvement in detection sensitivity and 5x improvement in resolution.",
        "[Technical Approach]": "Utilizes quantum entanglement and squeezed light states to enhance radar signal processing beyond classical limits.",
        "[Benefits]": "• 90% reduction in false positives\n• Real-time tracking of objects <1cm\n• 50% reduction in power consumption",

        # Add more mappings based on your template's placeholders
    }

    # Output path
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = Path(__file__).parent.parent.parent / "output" / f"filled_quad_chart_{timestamp}.pptx"

    # Create filler
    filler = QuadChartFiller(template_path)

    # Fill with data
    filler.fill_quad_chart(example_data)

    # Save
    saved_path = filler.save(output_path)
    print(f"Filled template saved to: {saved_path}")

if __name__ == "__main__":
    # Check if running with arguments
    if len(sys.argv) > 1:
        # Expected: script.py template_path data_json output_path
        if len(sys.argv) != 4:
            print("Usage: python quadChartFiller.py <template_path> <data_json> <output_path>")
            sys.exit(1)

        template = sys.argv[1]
        data = json.loads(sys.argv[2])
        output = sys.argv[3]

        result = fill_template_from_json(template, data, output)
        print(f"Success: {result}")
    else:
        # Run example
        main()