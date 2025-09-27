#!/usr/bin/env python3
"""
BD Bible PowerPoint Quad Chart Generator
Generates professional quad charts with GMU branding
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# GMU Brand Colors
GMU_GREEN = RGBColor(0, 102, 51)  # #006633
GMU_GOLD = RGBColor(255, 204, 51)  # #FFCC33
GMU_DARK = RGBColor(0, 85, 55)    # #005537
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)

class QuadChartGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Set slide size to standard (16:9)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def create_quad_chart(self, data):
        """Create a quad chart slide with the provided data"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add title
        self.add_title(slide, data.get('title', 'Quad Chart'))

        # Add client/subtitle
        self.add_subtitle(slide, data.get('client', ''))

        # Calculate quadrant dimensions
        margin = Inches(0.3)
        gap = Inches(0.15)

        # Accounting for title space
        top_offset = Inches(1.0)
        available_height = self.prs.slide_height - top_offset - (margin * 2)
        available_width = self.prs.slide_width - (margin * 2)

        quad_width = (available_width - gap) / 2
        quad_height = (available_height - gap) / 2

        # Create quadrants
        quadrants = [
            {
                'title': 'Technical Approach',
                'content': data.get('technical_approach', ''),
                'left': margin,
                'top': top_offset,
                'color': GMU_GREEN
            },
            {
                'title': 'Management Approach',
                'content': data.get('management_approach', ''),
                'left': margin + quad_width + gap,
                'top': top_offset,
                'color': GMU_GOLD
            },
            {
                'title': 'Past Performance',
                'content': data.get('past_performance', ''),
                'left': margin,
                'top': top_offset + quad_height + gap,
                'color': GMU_GOLD
            },
            {
                'title': 'Cost/Schedule',
                'content': data.get('cost_schedule', ''),
                'left': margin + quad_width + gap,
                'top': top_offset + quad_height + gap,
                'color': GMU_GREEN
            }
        ]

        for quad in quadrants:
            self.add_quadrant(
                slide,
                quad['title'],
                quad['content'],
                quad['left'],
                quad['top'],
                quad_width,
                quad_height,
                quad['color']
            )

        # Add logo placeholder if provided
        if data.get('logo_path'):
            self.add_logo(slide, data['logo_path'])

        # Add footer with date
        self.add_footer(slide, data.get('footer', ''))

        return slide

    def add_title(self, slide, title_text):
        """Add title to the slide"""
        title = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.1),
            self.prs.slide_width - Inches(0.6), Inches(0.5)
        )
        text_frame = title.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GMU_GREEN
        p.alignment = PP_ALIGN.CENTER

    def add_subtitle(self, slide, subtitle_text):
        """Add subtitle/client name"""
        if subtitle_text:
            subtitle = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.5),
                self.prs.slide_width - Inches(0.6), Inches(0.3)
            )
            text_frame = subtitle.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = subtitle_text
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER

    def add_quadrant(self, slide, title, content, left, top, width, height, header_color):
        """Add a quadrant with header and content"""
        # Create quadrant container (border)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        border.fill.solid()
        border.fill.fore_color.rgb = WHITE
        border.line.color.rgb = RGBColor(200, 200, 200)
        border.line.width = Pt(0.5)

        # Add header bar
        header_height = Inches(0.35)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, header_height
        )
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        header.line.color.rgb = header_color

        # Add header text
        header_text = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.05),
            width - Inches(0.2), header_height - Inches(0.1)
        )
        text_frame = header_text.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.02)
        text_frame.margin_bottom = Inches(0.02)

        p = text_frame.add_paragraph()
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

        # Add content area
        content_top = top + header_height + Inches(0.05)
        content_height = height - header_height - Inches(0.1)

        content_box = slide.shapes.add_textbox(
            left + Inches(0.1), content_top,
            width - Inches(0.2), content_height
        )

        text_frame = content_box.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True

        # Process content (support bullet points)
        if isinstance(content, str):
            lines = content.strip().split('\n') if content else []
        elif isinstance(content, list):
            lines = content
        else:
            lines = []

        for i, line in enumerate(lines):
            p = text_frame.add_paragraph()
            line = line.strip()

            # Check if line should be a bullet point
            if line.startswith('•') or line.startswith('-'):
                p.text = line[1:].strip()
                p.level = 0
                p.font.size = Pt(10)
            else:
                p.text = line
                p.font.size = Pt(10)

            p.font.name = 'Calibri'
            p.font.color.rgb = BLACK
            p.space_after = Pt(3)

            # Remove extra spacing from first paragraph
            if i == 0 and text_frame.paragraphs[0].text == '':
                text_frame.paragraphs[0].text = p.text
                text_frame.paragraphs[0].font.name = p.font.name
                text_frame.paragraphs[0].font.size = p.font.size
                text_frame.paragraphs[0].font.color.rgb = p.font.color.rgb
                text_frame.paragraphs.pop()

    def add_logo(self, slide, logo_path):
        """Add company logo to the slide"""
        if os.path.exists(logo_path):
            try:
                # Add logo in top-right corner
                slide.shapes.add_picture(
                    logo_path,
                    self.prs.slide_width - Inches(1.5),
                    Inches(0.1),
                    height=Inches(0.5)
                )
            except Exception as e:
                print(f"Warning: Could not add logo: {e}", file=sys.stderr)

    def add_footer(self, slide, footer_text=''):
        """Add footer with date and optional text"""
        footer_text = footer_text or f"Generated: {datetime.now().strftime('%Y-%m-%d')}"

        footer = slide.shapes.add_textbox(
            Inches(0.3),
            self.prs.slide_height - Inches(0.3),
            self.prs.slide_width - Inches(0.6),
            Inches(0.2)
        )

        text_frame = footer.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = footer_text
        p.font.name = 'Calibri'
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    def save(self, output_path):
        """Save the presentation to file"""
        self.prs.save(output_path)
        return output_path

def main():
    """Main function to handle command-line execution"""
    if len(sys.argv) < 2:
        print("Usage: python pptxGenerator.py <json_data> [output_path]", file=sys.stderr)
        sys.exit(1)

    try:
        # Parse input JSON
        json_data = json.loads(sys.argv[1])

        # Determine output path
        if len(sys.argv) > 2:
            output_path = sys.argv[2]
        else:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = f"/tmp/quad_chart_{timestamp}.pptx"

        # Generate the quad chart
        generator = QuadChartGenerator()
        generator.create_quad_chart(json_data)
        saved_path = generator.save(output_path)

        # Return the path as JSON
        result = {
            'success': True,
            'path': saved_path,
            'message': 'Quad chart generated successfully'
        }
        print(json.dumps(result))

    except json.JSONDecodeError as e:
        error = {
            'success': False,
            'error': f'Invalid JSON: {str(e)}'
        }
        print(json.dumps(error), file=sys.stderr)
        sys.exit(1)

    except Exception as e:
        error = {
            'success': False,
            'error': str(e)
        }
        print(json.dumps(error), file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()