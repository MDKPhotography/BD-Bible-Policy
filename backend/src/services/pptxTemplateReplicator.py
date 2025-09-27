#!/usr/bin/env python3
"""
PowerPoint Template Replicator
Preserves exact formatting while replacing placeholder text
"""

import json
import sys
import os
import re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
import traceback

class TemplateReplicator:
    def __init__(self, template_path):
        """
        Initialize with a PowerPoint template file
        Args:
            template_path: Path to the PPTX template file
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")

        self.template_path = template_path
        self.presentation = None
        self.placeholders_found = []

    def load_template(self):
        """Load the PowerPoint template"""
        self.presentation = Presentation(self.template_path)
        return self

    def find_placeholders(self):
        """
        Find all placeholder text in the template
        Looks for patterns like {{PLACEHOLDER_NAME}}
        """
        self.placeholders_found = []
        placeholder_pattern = re.compile(r'\{\{([A-Z_0-9]+)\}\}')

        for slide_idx, slide in enumerate(self.presentation.slides):
            # Check all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    matches = placeholder_pattern.findall(shape.text)
                    for match in matches:
                        self.placeholders_found.append({
                            'slide': slide_idx,
                            'shape': shape.name,
                            'placeholder': match,
                            'full_pattern': f'{{{{{match}}}}}'
                        })

                # Check table cells if shape is a table
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            matches = placeholder_pattern.findall(cell.text)
                            for match in matches:
                                self.placeholders_found.append({
                                    'slide': slide_idx,
                                    'shape': shape.name,
                                    'placeholder': match,
                                    'full_pattern': f'{{{{{match}}}}}',
                                    'is_table': True
                                })

        return self.placeholders_found

    def replace_text_preserve_format(self, shape, old_text, new_text):
        """
        Replace text while preserving all formatting
        This maintains font, size, color, bold, italic, etc.
        """
        if not hasattr(shape, 'text_frame'):
            return False

        text_frame = shape.text_frame

        # Process each paragraph
        for paragraph in text_frame.paragraphs:
            # Check if this paragraph contains our placeholder
            if old_text in paragraph.text:
                # Store original formatting from first run
                original_format = None
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    original_format = {
                        'font_name': first_run.font.name,
                        'font_size': first_run.font.size,
                        'bold': first_run.font.bold,
                        'italic': first_run.font.italic,
                        'underline': first_run.font.underline,
                        'color': first_run.font.color.rgb if first_run.font.color.type == 1 else None
                    }

                # Replace the text
                full_text = paragraph.text
                new_full_text = full_text.replace(old_text, new_text)

                # Clear runs and recreate with formatting
                paragraph.clear()
                run = paragraph.add_run()
                run.text = new_full_text

                # Restore formatting if we had it
                if original_format:
                    if original_format['font_name']:
                        run.font.name = original_format['font_name']
                    if original_format['font_size']:
                        run.font.size = original_format['font_size']
                    if original_format['bold'] is not None:
                        run.font.bold = original_format['bold']
                    if original_format['italic'] is not None:
                        run.font.italic = original_format['italic']
                    if original_format['underline'] is not None:
                        run.font.underline = original_format['underline']
                    if original_format['color']:
                        run.font.color.rgb = original_format['color']

                return True

        return False

    def replace_in_table(self, table, old_text, new_text):
        """Replace text in table cells while preserving formatting"""
        for row in table.rows:
            for cell in row.cells:
                if old_text in cell.text:
                    # Store the original formatting
                    text_frame = cell.text_frame
                    for paragraph in text_frame.paragraphs:
                        if old_text in paragraph.text:
                            # Preserve paragraph alignment
                            alignment = paragraph.alignment

                            # Store run formatting
                            original_format = None
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                original_format = {
                                    'font_name': first_run.font.name,
                                    'font_size': first_run.font.size,
                                    'bold': first_run.font.bold,
                                    'italic': first_run.font.italic,
                                    'color': first_run.font.color.rgb if first_run.font.color.type == 1 else None
                                }

                            # Replace text
                            new_paragraph_text = paragraph.text.replace(old_text, new_text)
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = new_paragraph_text

                            # Restore formatting
                            if original_format:
                                if original_format['font_name']:
                                    run.font.name = original_format['font_name']
                                if original_format['font_size']:
                                    run.font.size = original_format['font_size']
                                if original_format['bold'] is not None:
                                    run.font.bold = original_format['bold']
                                if original_format['italic'] is not None:
                                    run.font.italic = original_format['italic']
                                if original_format['color']:
                                    run.font.color.rgb = original_format['color']

                            # Restore alignment
                            paragraph.alignment = alignment

    def replace_placeholders(self, data_dict):
        """
        Replace all placeholders with actual data
        Args:
            data_dict: Dictionary mapping placeholder names to replacement values
                      e.g., {'CLIENT_NAME': 'Department of Defense', 'TECHNICAL': 'Advanced AI...'}
        """
        if not self.presentation:
            self.load_template()

        replacements_made = []

        for slide_idx, slide in enumerate(self.presentation.slides):
            # Process all shapes
            for shape in slide.shapes:
                # Handle text in shapes
                if hasattr(shape, "text"):
                    for placeholder, value in data_dict.items():
                        full_placeholder = f'{{{{{placeholder}}}}}'
                        if full_placeholder in shape.text:
                            # Handle multi-line values (convert \n to actual line breaks)
                            if isinstance(value, list):
                                value = '\n'.join(value)

                            success = self.replace_text_preserve_format(
                                shape, full_placeholder, str(value)
                            )
                            if success:
                                replacements_made.append({
                                    'slide': slide_idx,
                                    'placeholder': placeholder,
                                    'replaced_with': value
                                })

                # Handle tables
                if shape.has_table:
                    for placeholder, value in data_dict.items():
                        full_placeholder = f'{{{{{placeholder}}}}}'
                        # Check if placeholder exists in table
                        table_text = ''
                        for row in shape.table.rows:
                            for cell in row.cells:
                                table_text += cell.text

                        if full_placeholder in table_text:
                            if isinstance(value, list):
                                value = '\n'.join(value)

                            self.replace_in_table(shape.table, full_placeholder, str(value))
                            replacements_made.append({
                                'slide': slide_idx,
                                'placeholder': placeholder,
                                'replaced_with': value,
                                'in_table': True
                            })

        return replacements_made

    def process_bullets(self, text):
        """Convert bullet markers to proper bullet points"""
        lines = text.split('\n')
        processed = []
        for line in lines:
            line = line.strip()
            # Convert various bullet markers to standard format
            if line.startswith(('•', '-', '*', '○', '▪', '►')):
                line = '• ' + line[1:].strip()
            processed.append(line)
        return '\n'.join(processed)

    def save(self, output_path):
        """
        Save the modified presentation
        Args:
            output_path: Where to save the new PPTX file
        """
        if not self.presentation:
            raise ValueError("No presentation loaded")

        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        self.presentation.save(output_path)
        return output_path

    def get_template_info(self):
        """Get information about the template"""
        if not self.presentation:
            self.load_template()

        info = {
            'slide_count': len(self.presentation.slides),
            'slide_width': self.presentation.slide_width,
            'slide_height': self.presentation.slide_height,
            'placeholders': self.find_placeholders(),
            'slides': []
        }

        for idx, slide in enumerate(self.presentation.slides):
            slide_info = {
                'index': idx,
                'layout_name': slide.slide_layout.name,
                'shapes': []
            }

            for shape in slide.shapes:
                shape_info = {
                    'name': shape.name,
                    'type': shape.shape_type,
                    'has_text': hasattr(shape, 'text'),
                    'has_table': shape.has_table
                }
                if hasattr(shape, 'text'):
                    shape_info['text_preview'] = shape.text[:100] if shape.text else ''

                slide_info['shapes'].append(shape_info)

            info['slides'].append(slide_info)

        return info


def main():
    """Main function for command-line execution"""
    if len(sys.argv) < 3:
        print("Usage: python pptxTemplateReplicator.py <template_path> <json_data> [output_path]")
        sys.exit(1)

    template_path = sys.argv[1]

    try:
        # Parse JSON data
        json_data = json.loads(sys.argv[2])

        # Determine output path
        if len(sys.argv) > 3:
            output_path = sys.argv[3]
        else:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = f"/tmp/generated_{timestamp}.pptx"

        # Create replicator
        replicator = TemplateReplicator(template_path)
        replicator.load_template()

        # Find placeholders (for info)
        placeholders = replicator.find_placeholders()

        # Replace placeholders with data
        replacements = replicator.replace_placeholders(json_data)

        # Save the result
        saved_path = replicator.save(output_path)

        # Return result as JSON
        result = {
            'success': True,
            'path': saved_path,
            'placeholders_found': placeholders,
            'replacements_made': replacements,
            'message': f'Template processed successfully. {len(replacements)} replacements made.'
        }
        print(json.dumps(result))

    except FileNotFoundError as e:
        error = {
            'success': False,
            'error': f'Template file not found: {str(e)}'
        }
        print(json.dumps(error), file=sys.stderr)
        sys.exit(1)

    except json.JSONDecodeError as e:
        error = {
            'success': False,
            'error': f'Invalid JSON data: {str(e)}'
        }
        print(json.dumps(error), file=sys.stderr)
        sys.exit(1)

    except Exception as e:
        error = {
            'success': False,
            'error': str(e),
            'traceback': traceback.format_exc()
        }
        print(json.dumps(error), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()