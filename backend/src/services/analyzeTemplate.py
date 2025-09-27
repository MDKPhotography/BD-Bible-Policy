#!/usr/bin/env python3
"""
Analyze PowerPoint template to identify all elements, formatting, and structure
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

def analyze_font(font):
    """Extract font properties"""
    if not font:
        return None

    font_info = {}

    try:
        if font.name:
            font_info['name'] = font.name
        if font.size:
            font_info['size'] = font.size.pt if hasattr(font.size, 'pt') else font.size
        if font.bold is not None:
            font_info['bold'] = font.bold
        if font.italic is not None:
            font_info['italic'] = font.italic
        if font.underline is not None:
            font_info['underline'] = font.underline

        # Handle color carefully
        if hasattr(font, 'color') and font.color:
            try:
                if hasattr(font.color, 'rgb') and font.color.rgb:
                    font_info['color'] = str(font.color.rgb)
            except:
                pass
    except Exception as e:
        pass

    return font_info if font_info else None

def analyze_paragraph(paragraph):
    """Analyze paragraph formatting"""
    return {
        'text': paragraph.text,
        'alignment': str(paragraph.alignment) if paragraph.alignment else None,
        'level': paragraph.level,
        'font': analyze_font(paragraph.font),
        'runs': [
            {
                'text': run.text,
                'font': analyze_font(run.font)
            } for run in paragraph.runs
        ]
    }

def analyze_text_frame(text_frame):
    """Analyze text frame content and formatting"""
    if not text_frame:
        return None

    return {
        'text': text_frame.text,
        'paragraphs': [analyze_paragraph(p) for p in text_frame.paragraphs],
        'margin_left': text_frame.margin_left,
        'margin_right': text_frame.margin_right,
        'margin_top': text_frame.margin_top,
        'margin_bottom': text_frame.margin_bottom,
        'word_wrap': text_frame.word_wrap
    }

def analyze_shape(shape, index):
    """Analyze individual shape properties"""
    shape_info = {
        'index': index,
        'name': shape.name,
        'shape_type': str(shape.shape_type),
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height
    }

    # Check for text
    if shape.has_text_frame:
        shape_info['text_frame'] = analyze_text_frame(shape.text_frame)

    # Check for fill - with better error handling
    try:
        if hasattr(shape, 'fill') and shape.fill:
            fill = shape.fill
            if fill.type is not None:
                shape_info['fill'] = {'type': str(fill.type)}
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    try:
                        if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                            shape_info['fill']['fore_color'] = str(fill.fore_color.rgb)
                    except:
                        pass
    except Exception as e:
        pass

    # Check for line - with better error handling
    try:
        if hasattr(shape, 'line') and shape.line:
            line = shape.line
            shape_info['line'] = {}
            if hasattr(line, 'width'):
                shape_info['line']['width'] = line.width
            if hasattr(line, 'color') and line.color:
                try:
                    if hasattr(line.color, 'rgb') and line.color.rgb:
                        shape_info['line']['color'] = str(line.color.rgb)
                except:
                    pass
    except Exception as e:
        pass

    # Check if it's a picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        shape_info['is_picture'] = True
        shape_info['image_format'] = shape.image.ext if hasattr(shape, 'image') else None

    # Check if it's a table
    if shape.has_table:
        table = shape.table
        shape_info['is_table'] = True
        shape_info['table'] = {
            'rows': len(table.rows),
            'columns': len(table.columns),
            'cells': []
        }
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_info = {
                    'row': row_idx,
                    'col': col_idx,
                    'text': cell.text,
                    'text_frame': analyze_text_frame(cell.text_frame) if cell.text_frame else None
                }
                shape_info['table']['cells'].append(cell_info)

    # Check if it's a group
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        shape_info['is_group'] = True
        shape_info['shapes'] = []
        for idx, sub_shape in enumerate(shape.shapes):
            shape_info['shapes'].append(analyze_shape(sub_shape, idx))

    return shape_info

def analyze_slide(slide, slide_index):
    """Analyze all elements in a slide"""
    slide_info = {
        'slide_number': slide_index + 1,
        'layout_name': slide.slide_layout.name if slide.slide_layout else None,
        'shapes': []
    }

    # Analyze all shapes
    for idx, shape in enumerate(slide.shapes):
        try:
            shape_analysis = analyze_shape(shape, idx)
            slide_info['shapes'].append(shape_analysis)
        except Exception as e:
            print(f"Warning: Could not fully analyze shape {idx}: {e}")
            slide_info['shapes'].append({
                'index': idx,
                'name': shape.name if hasattr(shape, 'name') else 'Unknown',
                'error': str(e)
            })

    # Get placeholder information
    slide_info['placeholders'] = []
    for shape in slide.placeholders:
        placeholder_info = {
            'idx': shape.placeholder_format.idx,
            'type': str(shape.placeholder_format.type) if shape.placeholder_format.type else None,
            'name': shape.name,
            'text': shape.text if shape.has_text_frame else None
        }
        slide_info['placeholders'].append(placeholder_info)

    return slide_info

def analyze_template(template_path):
    """Main function to analyze the PowerPoint template"""
    print(f"\n{'='*60}")
    print(f"POWERPOINT TEMPLATE ANALYSIS")
    print(f"Template: {template_path}")
    print(f"{'='*60}\n")

    # Load presentation
    prs = Presentation(template_path)

    # Basic presentation info
    print(f"PRESENTATION INFO:")
    print(f"  - Total slides: {len(prs.slides)}")
    print(f"  - Slide width: {prs.slide_width} EMUs ({prs.slide_width/914400:.2f} inches)")
    print(f"  - Slide height: {prs.slide_height} EMUs ({prs.slide_height/914400:.2f} inches)")

    # Analyze each slide
    analysis = {
        'template_file': str(template_path),
        'slide_count': len(prs.slides),
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height,
        'slides': []
    }

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'-'*40}")
        print(f"SLIDE {slide_idx + 1}:")
        print(f"{'-'*40}")

        slide_analysis = analyze_slide(slide, slide_idx)
        analysis['slides'].append(slide_analysis)

        # Print summary
        print(f"  Layout: {slide_analysis['layout_name']}")
        print(f"  Total shapes: {len(slide_analysis['shapes'])}")
        print(f"  Placeholders: {len(slide_analysis['placeholders'])}")

        # List text content
        print(f"\n  TEXT CONTENT:")
        for shape in slide_analysis['shapes']:
            if 'text_frame' in shape and shape['text_frame']:
                text = shape['text_frame']['text'].strip()
                if text:
                    print(f"    [{shape['index']}] {shape['name']}: \"{text[:50]}{'...' if len(text) > 50 else ''}\"")

        # List images
        images = [s for s in slide_analysis['shapes'] if s.get('is_picture')]
        if images:
            print(f"\n  IMAGES:")
            for img in images:
                print(f"    [{img['index']}] {img['name']} - Position: ({img['left']}, {img['top']})")

        # List tables
        tables = [s for s in slide_analysis['shapes'] if s.get('is_table')]
        if tables:
            print(f"\n  TABLES:")
            for tbl in tables:
                print(f"    [{tbl['index']}] {tbl['name']} - {tbl['table']['rows']}x{tbl['table']['columns']}")

    # Save detailed analysis to JSON
    output_path = Path(template_path).parent / 'template_analysis.json'
    with open(output_path, 'w') as f:
        json.dump(analysis, f, indent=2, default=str)

    print(f"\n{'='*60}")
    print(f"ANALYSIS COMPLETE")
    print(f"Detailed analysis saved to: {output_path}")
    print(f"{'='*60}\n")

    return analysis

if __name__ == "__main__":
    # Get template path
    template_path = Path(__file__).parent.parent.parent / "templates" / "powerpoint" / "space-bd-quad-charts.pptx"

    if not template_path.exists():
        print(f"Error: Template not found at {template_path}")
        sys.exit(1)

    # Run analysis
    analyze_template(template_path)